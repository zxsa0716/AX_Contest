"""Generate full 50-slide PPTX with all 17 figures embedded.

Self-contained alternative to Gamma (no API credits needed).
Uses python-pptx with custom layouts.
Output: report/AX_contest_2026_presentation.pptx
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[2]
FIGS = ROOT / "figs"
OUT = ROOT / "report" / "AX_contest_2026_presentation.pptx"

# Brand colors
PRIMARY = RGBColor(0x1D, 0x4E, 0xD8)  # blue
SECONDARY = RGBColor(0x15, 0x80, 0x3D)  # green
WARNING = RGBColor(0xD5, 0x5E, 0x00)  # orange/red
DARK = RGBColor(0x1A, 0x1A, 0x18)
GRAY = RGBColor(0x60, 0x60, 0x60)


def _add_title_slide(prs, title, subtitle, footer=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    bg.line.fill.background()
    # Top color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Emu(150_000))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    # Title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = DARK
    # Subtitle
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = "\n" + subtitle
        r2.font.size = Pt(18)
        r2.font.color.rgb = GRAY
    # Footer
    if footer:
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5))
        ftf = ft.text_frame
        fp = ftf.paragraphs[0]
        fp.alignment = PP_ALIGN.CENTER
        fr = fp.add_run()
        fr.text = footer
        fr.font.size = Pt(11)
        fr.font.color.rgb = GRAY


def _add_content_slide(prs, title, body_text, image_path=None, body_size=14):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = PRIMARY
    # Underline
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.0), Inches(12.5), Emu(20_000))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()

    # Body + Image layout
    if image_path and image_path.exists():
        # Image right, text left
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5))
        try:
            slide.shapes.add_picture(str(image_path), Inches(6.3), Inches(1.3),
                                     width=Inches(6.7), height=Inches(5.4))
        except Exception:
            pass
    else:
        # Body full width
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.4), Inches(5.6))

    btf = body_box.text_frame
    btf.word_wrap = True
    for line_text in body_text.strip().split("\n"):
        p = btf.add_paragraph()
        if line_text.startswith("• ") or line_text.startswith("- "):
            r = p.add_run()
            r.text = line_text[2:]
            r.font.size = Pt(body_size)
            r.font.color.rgb = DARK
            p.level = 1
        elif line_text.startswith("**") and line_text.endswith("**"):
            r = p.add_run()
            r.text = line_text.strip("*")
            r.font.size = Pt(body_size + 2)
            r.font.bold = True
            r.font.color.rgb = PRIMARY
        else:
            r = p.add_run()
            r.text = line_text
            r.font.size = Pt(body_size)
            r.font.color.rgb = DARK

    # Slide number footer
    foot = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3))
    fp = foot.text_frame.paragraphs[0]
    fp.alignment = PP_ALIGN.RIGHT
    fr = fp.add_run()
    fr.text = f"AX 2026 · {len(prs.slides)}"
    fr.font.size = Pt(9)
    fr.font.color.rgb = GRAY


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)

    # ============ DECK 1: 도입 + 문제의식 (10) ============
    _add_title_slide(prs,
        "한국 코스피 상장기업\n온실가스 공시 신뢰성 3중 검증",
        "GIR × ESG × Sentinel-5P × ODIAC 4중 비교와\nKSSB 2028 의무공시 검증체계 설계",
        "2026 AX 아이디어 경진대회 자유분석 부문 · 마감 2026-05-18")

    _add_content_slide(prs,
        "Hero Finding — 한 줄 정리",
        """**포스코홀딩스 Mann-Kendall τ**

GIR 법정 공시: τ = +1.00 (완벽 상승)
ESG 자체보고: τ = +0.67 (상승)
위성 NO₂:    τ = −1.00 (완벽 하강)
ODIAC CO₂:   τ = −1.00 (완벽 하강)

→ 한국 코스피 상장기업 중 GIR 공시와
  위성 독립측정의 방향이 완벽히 반대인
  사례를 처음으로 정량적으로 식별

→ 2028년 KSSB 의무공시 시행 직전
  검증체계 설계의 골든 타임""",
        FIGS / "fig_concept_4channel.png")

    _add_content_slide(prs,
        "한국 ESG 공시의 두 채널 — 구조적 불일치 가능성",
        """**채널 1: GIR 법정 명세서**
- 환경부 온실가스종합정보센터 신고
- 사업장 단위, 처벌 구속력 (배출권거래법)
- K-ETS 할당의 직접 근거

**채널 2: ESG 자체보고**
- GRI 305-1, TCFD, ISSB 기준
- 연결 기준, 투자자 대상 자율 공시
- 허위 공시 실질 제재 부재

**불일치 발생 4가지 원인**
- 조직경계 (control vs equity)
- 배출계수 (국내고시 vs 국제기준)
- 공간범위 (국내 vs 해외 포함)
- 처벌 비대칭 (제재 유인 차이)""")

    _add_content_slide(prs,
        "검증 시스템 부재 — 제3자 검증의 한계",
        """**현행 제3자 검증 (ISAE 3410, AA1000AS)**
- 자료의 적정성 확인에 그침
- 절대값 정확성 검증 X
- 독립적 물리 측정 기반 부재

**KSSB 2028 외부 검증 요건화**
- 2026-02-26 KSSB 제2호 확정
- FY27부터 의무공시 (2028 보고)
- 검증 프로토콜 미정 → 본 연구의 진입점

**위성 + ODIAC 독립 검증의 가능성**
- Sentinel-5P TROPOMI (Kim 2020 R=0.96)
- ODIAC 1km top-down CO₂ (Ahn-Goldberg 2025)
- ERA5 기상보정 표준 (Fioletov 2025)""")

    _add_content_slide(prs,
        "본 연구의 4중 비교 프레임",
        """**4채널 통합 비교**

GIR (정부) ←→ ESG (시장)
        ↓
위성 4종 (대기 물리)
+ ODIAC CO₂ (top-down)
        ↓
ERA5 기상보정 잔차
        ↓
Mann-Kendall 방향 일관성
        ↓
패턴 5종 분류 (A·B·C·D·E)
        ↓
검증 우선순위 매트릭스
→ KEITI · KSSB 정책 직결""",
        FIGS / "fig_concept_4channel.png")

    _add_content_slide(prs,
        "KSSB 2028 임박성 — 검증 체계 설계의 골든 타임",
        """**확정된 정책 타임라인**
- 2026-02-25: FSC 'ESG 공시 로드맵(안)' 발표
- 2026-02-26: KSSB 제1호·제2호·제101호 확정
- 2027 회계연도: 첫 의무공시 적용 (KOSPI 30조원↑)
- 2028년: 첫 보고서 제출 (약 49-58개사)

**골든 타임의 의미**
2026~2027년 2년이 검증 체계를 설계할
마지막 기회. 의무화는 형식적 제출 요건에
그칠 위험 → 본 연구가 즉시 적용 가능한
검증 프레임워크 제공.

**Gold 23개사 = KSSB 1차 대상과 직접 교집합**""")

    _add_content_slide(prs,
        "Gold 23개사 — KSSB 2028 ∩ KOSPI ∩ GIR ≥3yr",
        """**산업 구성 (5개 군)**

산업·에너지 (5):
  POSCO·현대제철·KEPCO·SK이노·LG엔솔

반도체·전자 (3):
  삼성전자·SK하이닉스·LG디스플레이

석유화학 (3):
  롯데케미칼·한화솔루션·한화

금융·서비스 (4):
  삼성생명·IBK·KT·네이버

지주·기타 (8):
  삼성물산·CJ제일제당·롯데쇼핑·이마트·
  두산·대한항공·현대차·현대모비스""",
        FIGS / "fig_map_gold_sites.png")

    _add_content_slide(prs,
        "데이터 인프라 — 18 데이터셋 통합",
        """**위성·기상 (5)**
- Sentinel-5P NO₂·SO₂·CO·HCHO (1380 monthly)
- ODIAC v2024 CO₂ 1km (60 monthly rasters)
- ERA5-Land + ERA5 BLH (5 변수)
- MERRA-2 sensitivity check
- ASOS 5 stations × 5 years

**공식 공시 (8)**
- GIR 명세서 7년 (7,998행)
- 지속가능경영보고서 126 PDFs (1.5GB)
- DART 사업보고서 + 재무
- K-ETS 할당 4 phase
- KCGS ESG 등급 21건
- 통합환경허가 1,065건
- NIR 국가 인벤토리
- DART 자율공시 81건 ESG

**기준·매핑 (5)**
- KOSPI 전체 789 corp_code
- KSSB 2028 pool 49사
- VWorld 지오코딩
- KCGS 21 supervised label
- 4 ADR 의사결정""",
        FIGS / "fig_map_odiac_korea.png")

    _add_content_slide(prs,
        "Agenda — 5-Part 발표 구성 (50 슬라이드)",
        """**Part 1: 도입 + 문제의식 (1-10)**
KSSB 2028 임박성, 4중 비교 프레임

**Part 2: 데이터 + 위성 영상 (11-20)**
18 데이터셋 명세, 한국 위성 지도

**Part 3: 분석 8단계 방법론 (21-30)**
ERA5 보정, 패턴 분류, Heckman

**Part 4: 결과 — 결정적 발견 4가지 (31-40)**
패턴 D 2개사, 이상탐지, 회귀, SHAP

**Part 5: 정책 + 결론 (41-50)**
3 정책 카드, KSSB 활용, 한계 + 향후""")

    _add_content_slide(prs,
        "Part 1 결론: 4중 비교 프레임의 신규성",
        """**선행 연구 격차**
- Liu 2020 (Nature): 중국 발전·산업 시설
- Kim 2020 (Atmosphere): 한국 TROPOMI vs CAPSS
- Fioletov 2025 (ACP): ERA5 보정 표준화
- Ahn-Goldberg 2025 (AGU Adv): 54 cities CO₂
- Taean 2024 (ScienceDirect): 한국 SO₂ top-down

**본 연구 차별점**
- 한국 코스피 기업 단위 4중 비교 최초
- KSSB 2028 정책 직결 + Gold 직접 매칭
- ERA5+MERRA-2+ASOS 3-layer 기상 보정
- 자동화 수집·파싱 시스템 영구 편입""")

    # ============ DECK 2: 데이터 + 위성 영상 (11-20) ============
    _add_title_slide(prs,
        "Part 2 — 데이터 인프라",
        "18 데이터셋 통합 + 한국 위성 영상")

    _add_content_slide(prs,
        "한국 ODIAC CO₂ 전국 지도 (2023-05)",
        """**시각화 핵심**
- ODIAC v2024 1km 해상도
- 한국 영토 클립 (124°E-132°E, 33°N-39°N)
- log scale color (hot palette)
- Gold 23개사 사업장 ▲ overlay

**관찰**
- 수도권 + 동남부 산업벨트 집중
- 포항·울산·여수·당진 hotspot
- 한국 연간 ~180 MtC, 월간 ~16 MtC

**데이터 소스**
db.cger.nies.go.jp/dataset/ODIAC2024
60 monthly GeoTIFF (2019-2023)""",
        FIGS / "fig_map_odiac_korea.png")

    _add_content_slide(prs,
        "ODIAC 4계절 비교 — 산업 활동 패턴",
        """**4 panel: 1월·4월·7월·10월 (2023)**

겨울 (1월):
  난방 부하 + 화력 발전 가동률 상승

봄 (4월):
  생산 활동 정점, 안정적 패턴

여름 (7월):
  냉방 부하 + 일부 정유 가동률

가을 (10월):
  화학·정유 정기보수 영향

→ 산업단지의 계절성이 위성에서 가시화""",
        FIGS / "fig_map_odiac_seasonal.png")

    _add_content_slide(prs,
        "Gold 23개사 사업장 위치 (산업별)",
        """**5개 산업 분류 색상**
- 철강 (steel): 빨강
- 석유화학 (petrochem): 주황
- 발전 (power_coal): 검정
- 반도체 (semiconductor): 하늘색
- 금융·기타 (other): 회색

**사업장 좌표 출처**
- VWorld Geocoder API 2.0
- 23/23 매칭 성공 (100%)
- 단, 일부 firm은 본사 좌표 (한계 8.0 참조)

**클러스터 패턴**
- 수도권: 본사·서비스 firms 집중
- 동남부: 포항·울산 산업밸트
- 남해안: 여수·광양 석유화학""",
        FIGS / "fig_map_gold_sites.png")

    _add_content_slide(prs,
        "ASOS 지상 관측 + 23사 매칭",
        """**ASOS 5개 지점 활용**
- 서울 (108): 수도권 firms
- 수원 (119): 분당·기흥
- 인천 (112): 인천 공장
- 포항 (138): POSCO 포항
- 광주 (156): 호남권 firms

**평균 매칭 거리**
- 9.2 km (대부분 < 20 km)
- 일부 site > 20 km (당진·태안 flag)

**ERA5 vs ASOS 검증 목적**
- 격자 모델 vs 지점 관측 sanity
- 기상보정 robustness 추가 신뢰도""",
        FIGS / "fig_map_asos_stations.png")

    _add_content_slide(prs,
        "데이터 커버리지 매트릭스",
        """**비율 요약 (157 firm-year 기준)**

GIR Scope 1:        100% (157/157)
ESG Scope 1:         76% (119/157)
Sentinel-5P NO₂:    100% (115/115 in 2019-23)
Sentinel-5P SO₂:     96% (110/115)
Sentinel-5P CO:     100%
Sentinel-5P HCHO:   100%
ERA5-Land:          100%
ERA5 BLH:           100%
MERRA-2 PBLTOP:     100%
ODIAC CO₂ 1km:      100% (1380/1380 monthly)
ASOS 5 stations:    100% (1380 monthly)

**ESG 76% — 잔여 24% 사유**
- 11 구조적 미발간 (LG엔솔 pre-2022,
  현대차 pre-2022, 네이버 2019, IBK 2020,
  롯데쇼핑·이마트 2019-2020 등)""",
        FIGS / "fig_gir_heatmap.png")

    _add_content_slide(prs,
        "GIR Scope 1 5년 추이 (Gold 23사)",
        """**시계열 관찰 (2019-2023)**

대부분 firms 일관 하강 추세
→ 탈탄소 정책 + 코로나 영향

**Top 6 emitters (절대값)**
1. 한국전력공사 (압도적 1위)
2. POSCO홀딩스
3. 현대제철
4. 한화솔루션
5. SK이노베이션
6. 삼성전자

**산업별 패턴**
- 발전 (KEPCO): 단계적 감축 진행
- 철강 (POSCO·현대제철): 변동
- 석유화학: 연도별 상이
- 반도체: 일정 (생산 확대 + 효율 개선)""",
        FIGS / "fig_gir_timeseries.png")

    _add_content_slide(prs,
        "ESG 보고서 자동 수집 시스템",
        """**3-Source Fallback 파이프라인**
1. DART Open API 자율공시 → URL
2. KRX ESG 포털 (Selenium)
3. IR 사이트 직접 다운로드

**수집 결과 (126 PDFs, 1.5GB)**
- 23 firms × 5-7 년치 (2019-2025 일부)
- HIGH 신뢰도 21 / MED 91 / LOW 14
- 99% Scope 1 추출 성공

**파싱 항목 (GRI 305-1/2/3)**
- Scope 1 (직접 배출)
- Scope 2 location vs market
- Scope 3 카테고리
- 조직경계 (operational/financial/equity)
- 제3자 검증 (ISAE 3410 / AA1000AS)
- 보고기준 (GRI / IFRS S2 / TCFD / KSSB)

**자동화 시스템 영구 편입**
.claude/skills/sustainability-report-collect.md
.claude/skills/esg-scope-extract.md""")

    _add_content_slide(prs,
        "위성 + ERA5 통합 패널 구조",
        """**Panel: 23 sites × 60 months = 1380 rows**

각 row 변수:
- company_id, site_id, year, month
- no2_mean, so2_mean, co_mean, hcho_mean
- no2_resid, so2_resid, co_resid, hcho_resid (ERA5 보정 잔차)
- era5_u10, v10, t2m, tp, blh
- merra2_pbltop, ps, disph, qv2m
- odiac_sum_tC (월 1km buffer)
- asos_avgTa, avgWs, sumRn

**ERA5 기상보정 R² (Per-site OLS)**
- HCHO: 0.94 (계절성 강함)
- SO₂:  0.79
- NO₂:  0.76
- CO:   0.67

→ 잔차 = 배출 활동 신호 (기상 효과 제거)""",
        FIGS / "fig_satellite_scatter.png")

    _add_content_slide(prs,
        "Part 2 결론: 데이터 인프라의 완결성",
        """**18 데이터셋 통합의 의의**
- 단일 firm 단일 비교가 아닌 패널 분석 가능
- 기상 변동성 통제 (ERA5+MERRA-2+ASOS 3-layer)
- 위성 4종 + ODIAC = 5 채널 독립 측정
- 자동화 시스템 영구 편입 (재현 가능)

**한국 코스피 ESG 분석 인프라 구축의 첫 걸음**
- GitHub 공개: https://github.com/zxsa0716/AX_Contest
- API 키 발급 → 6-8시간 재현 가능

**다음 Part 3: 8단계 분석 방법론**""")

    # ============ DECK 3: 분석 방법론 (21-30) ============
    _add_title_slide(prs,
        "Part 3 — 8단계 분석 방법론",
        "ERA5 보정 → 패턴 분류 → 이상탐지 → Heckman → SHAP")

    _add_content_slide(prs,
        "8단계 파이프라인 전체 도식",
        """**1. 데이터 수집** (수동 + 자동화)
**2. 전처리** (7-step: 매칭·Tier·MICE·신뢰도 등)
**3. ERA5 기상보정** (per-site OLS, R² 0.67-0.94)
**4. 괴리 지표** (절대·상대·방향)
**5. 이상탐지 3층** (IF + LOF + MK + KCGS)
**6. 4중 비교 + 패턴 5종** (Mann-Kendall τ)
**7. Heckman 2-stage 회귀** (+ Bootstrap CI)
**8. SHAP TreeExplainer** + 정책 매트릭스

각 단계별 출력 → data/processed/*.csv
재현 명령: `python src/run_all_analysis.py`""",
        FIGS / "fig_concept_4channel.png")

    _add_content_slide(prs,
        "단계 3: ERA5 기상보정 (Fioletov 2025 ACP)",
        """**OLS 회귀 (per-site, monthly)**

NO₂_raw = β₀ + β₁·u10 + β₂·v10 + β₃·t2m
        + β₄·tp + β₅·blh + month_dummies + ε

**잔차 ε = 배출 활동 신호**

**R² 분포 (per-site mean)**
- HCHO: 0.94 (계절성 매우 강함)
- SO₂:  0.79
- NO₂:  0.76
- CO:   0.67

→ 원시 위성 신호의 67-94%가 기상 변동
→ 보정 후 잔차로 패턴 분석

**Robustness check**
- ERA5 vs MERRA-2 sensitivity
- ASOS 지상 관측으로 ERA5 검증
- 보정 전후 패턴 분류 동일 확인""",
        FIGS / "fig_satellite_scatter.png")

    _add_content_slide(prs,
        "단계 5: 이상탐지 3층 앙상블",
        """**Layer 1: 횡단면 (Cross-sectional)**
- Isolation Forest (n_estimators=200)
- Local Outlier Factor (k=20)
- contamination sweep 0.05-0.20
- 두 모델 동시 이상 = anomaly

**Layer 2: 시계열 (Longitudinal)**
- Mann-Kendall τ + p<0.1
- |τ| ≥ 0.4 + 유의 = trend anomaly

**Layer 3: 지도학습 (Supervised)**
- KCGS 2023-2025 분기 등급조정 21건
- 등급 강등 = 외부 검증 라벨
- precision/recall으로 모델 calibration

**최종 분류**
- structural (L1∩L2): 4건 (KEPCO 2020-23)
- transient (L1만): 4건 (POSCO 3yr+SK)
- longitudinal (L2만): 14건
- normal: 93건""",
        FIGS / "fig_top6_multipanel.png")

    _add_content_slide(prs,
        "단계 6: Mann-Kendall 4중 비교 + 패턴 5종",
        """**4 시계열 × 23 firms (Gold)**
- GIR Scope 1
- ESG Scope 1
- 위성 NO₂ (ERA5 보정 잔차)
- ODIAC CO₂

**각각 τ 계산 → 부호 조합으로 분류**

**패턴 정의**
- A_up:    4 채널 모두 ↑
- A_down:  4 채널 모두 ↓
- B:       ESG만 반대 (B_esg_suspect)
- C:       GIR만 반대 (C_gir_suspect)
- D:       2+ 채널 반대 (최심각)
- E:       무추세 / mixed

**Threshold (소샘플 N=5)**
- |τ| ≥ 0.4 (방향 판정)
- p < 0.1 (유의성)""",
        FIGS / "fig_mk_tau_forest.png")

    _add_content_slide(prs,
        "단계 7: Heckman 2-stage + Bootstrap CI",
        """**Stage 1 (Selection): Probit**
P(ESG_보고 발간 | 기업 특성) = Φ(...)

→ Inverse Mills Ratio (IMR) 계산

**Stage 2 (Outcome): FE Panel OLS**
disc_pct = α + β₁·ln(GIR) + β₂·in_KSSB
         + β₃·IMR + 산업·연도 더미
         + 클러스터 SE (firm)

**Bootstrap 95% CI (B=2000, firm-block)**
- 소샘플 (N=104) 불확실성 정량화
- ln(GIR) -2.00 [-8.93, -0.21]
- yr_2021 +11.49 [0.13, 32.98]

**해석**
대기업일수록 괴리율 낮음
2021년 일시적 괴리 급증 (코로나 회복기)""")

    _add_content_slide(prs,
        "단계 8: SHAP TreeExplainer + 정책 매트릭스",
        """**SHAP feature_perturbation='interventional'**
- Path-dependent 편향 회피
- 8 features: log_GIR, NO₂, SO₂, CO, HCHO,
  ODIAC, ERA5_BLH, energy_TJ

**Feature Importance (mean |SHAP|)**
1. NO₂_resid (위성 보정값)
2. log(GIR Scope 1)
3. 괴리율 DIFF_rel
4. SO₂_resid
5. ODIAC CO₂
6. CO_resid
7. ERA5 BLH
8. energy_TJ

**검증 우선순위 매트릭스**
priority_score = 0.4·괴리 + 0.4·위성불일치
              + 0.2·이상등급
- 상위 25% → 즉시 검증
- 26-50% → 우선 관찰
- 51-100% → 일반 모니터링""",
        FIGS / "fig_shap_summary.png")

    _add_content_slide(prs,
        "괴리 지표 설계 + 산업별 분포",
        """**3가지 괴리 지표**
- 절대 괴리: ESG_S1 - GIR_S1 (tCO₂eq)
- 상대 괴리율: 100·(ESG-GIR)/GIR (%)
- 방향성: sign(괴리율)

**6개 분석 차원**
- 규모 분포 (P10/P50/P90)
- 방향성 비율 (과소·과대 신고)
- 연도별 추세 (2019→2023)
- 업종별 비교 (energy/steel/chem/finance)
- 검증 여부 효과 (3rd party assurance)
- GIR Tier 효과 (T1/T2/T3)

**산업별 패턴**
- 철강·발전: 대규모 + 안정적
- 반도체: 일정 + 효율 개선
- 화학: 연도별 변동
- 금융·서비스: Scope 1 미미""",
        FIGS / "fig_gir_heatmap.png")

    _add_content_slide(prs,
        "자동화 수집·파싱 시스템 (영구 deliverable)",
        """**구조 (.claude/skills/)**
1. sustainability-report-collect.md
2. esg-scope-extract.md

**수집 파이프라인**
sustainability_report_collector.py
- DART Open API → 자율공시 검색
- KRX ESG 포털 → Selenium fallback
- IR 사이트 → URL pattern matching
- SHA-256 deduplication
- Resume-safe (interrupted 시 재개)

**파싱 파이프라인**
sustainability_report_parser.py
- pdfplumber + 정규식
- GRI 305-1 자동 추출
- 신뢰도 플래그 (HIGH/MED/LOW)
- 단위 정규화 (tCO₂eq)
- 조직경계·검증·보고기준 추출

**재현성**
GitHub public + 4 ADR + .env 템플릿
→ 누구나 6-8시간 만에 재현 가능""")

    _add_content_slide(prs,
        "Part 3 결론: 방법론적 강건성",
        """**8단계 통합 파이프라인**
- 모든 단계가 학술적 근거 보유
- 재현 가능 (run_all_analysis.py 1 명령)
- 4 ADR로 의사결정 추적

**Robustness Checks**
- ERA5 vs MERRA-2 sensitivity
- contamination sweep (0.05-0.20)
- Bootstrap CI (B=2000, firm-block)
- ASOS 지상 검증
- 패턴 ERA5 보정 전후 동일

**다음 Part 4: 결정적 발견 4가지**""")

    # ============ DECK 4: 결과 (31-40) ============
    _add_title_slide(prs,
        "Part 4 — 결정적 발견 4가지",
        "패턴 D 2개사 / 이상탐지 / 회귀 / SHAP")

    _add_content_slide(prs,
        "발견 ① 패턴 5종 분포 (Gold 23사)",
        """**Mann-Kendall 4중 비교 결과**

A_consistent_down: 12개사 (52%)
  탈탄소 정책 정합 — 두산·한화·KEPCO·KT·LGD·현대차 등

mixed: 7개사 (30%)
  부분 일치 — SK하이닉스·현대제철·롯데쇼핑·한화솔루션 등

**D_both_suspect: 2개사 (9%) ⚠️**
  포스코홀딩스·삼성전자 — 공시↑ but 위성·ODIAC↓

C_gir_suspect: 1개사 (4%)
  현대모비스 — GIR↓·ESG↑ + 위성↓

A_consistent_up: 1개사 (4%)
  네이버 — 데이터센터 확장 정합""",
        FIGS / "fig_pattern_distribution.png")

    _add_content_slide(prs,
        "발견 ② 패턴 D 사례 ① — 포스코홀딩스",
        """**Mann-Kendall τ (2019-2023)**
- GIR Scope 1: τ = +1.00 (완벽 상승)
- ESG Scope 1: τ = +0.67 (상승)
- 위성 NO₂:    τ = −1.00 (완벽 하강)
- ODIAC CO₂:   τ = −1.00 (완벽 하강)

**해석 (3가지 가설)**
H1. 보고경계 확대 (해외 사업장 통합)
H2. Scope 정의 변경 (간접 → 직접 분류)
H3. 실제 배출 감소 + 보고 시차

**본 분석으로 단정 불가**
- 인과 추론 시도 X
- 기술적 발견: "방향 불일치 관찰"
- 후속 검증 필요 (KSSB 시행 시 우선순위)""",
        FIGS / "fig_case_studies.png")

    _add_content_slide(prs,
        "발견 ② 패턴 D 사례 ② — 삼성전자",
        """**Mann-Kendall τ (2019-2023)**
- GIR Scope 1: τ = +0.60
- ESG Scope 1: τ = +1.00 (최대 상승)
- 위성 NO₂:    τ = −0.40
- ODIAC CO₂:   τ = −0.40

**ESG 자체보고가 가장 강한 상승 추세**
- 사업 확장 (반도체 fab 증설)에 따른
  공시 절대값 증가 보고
- 그러나 수원·화성 buffer NO₂는 하강
  (효율 개선 또는 측정 한계)

**해석 caveat**
- 수원 본사 좌표 (HQ + 화성 fab 일부)
- buffer 10km 내 도시 배경 영향 가능
- 그럼에도 5년 일관 방향 차이는 유의""",
        FIGS / "fig_industrial_no2_timeseries.png")

    _add_content_slide(prs,
        "패턴 한국 지도 + 산업 맥락",
        """**지도 시각화**
- ★ 빨강: D 패턴 (포스코·삼성전자)
- ◆ 주황: C 패턴 (현대모비스)
- ● 초록: A up (네이버)
- ● 파랑: A down (12개사)
- ● 회색: mixed (7개사)

**지리적 패턴**
- 동남부 산업밸트 D 집중 (포스코)
- 수도권 mixed/down 다수
- 호남 (KEPCO) structural

**Implication**
KSSB 49사 의무공시 1차 적용 시
패턴 D·C 기업을 우선 검증 대상으로
지정 → 정책 자원 효율 배분""",
        FIGS / "fig_map_patterns.png")

    _add_content_slide(prs,
        "발견 ③ 이상탐지 3층 앙상블 (115 firm-year)",
        """**최종 분류**

structural (L1∩L2): 4건
  → 한국전력공사 2020-2023 (4년 연속)
  → 단일 최대 배출자 절대값 효과

**transient (L1만): 4건**
  → 포스코홀딩스 2021·2022·2023 (3년 연속) ⚠️
  → SK하이닉스 2021

longitudinal (L2만): 14건
  → 5년 추세는 비정상이나 단년 정상
  → 추세 안정성 의심 firms

normal: 93건 (81%)
  → 두 layer 모두 정상

**Contamination Sensitivity**
- c=0.05 → 4 anomalies
- c=0.10 → 8 anomalies (default)
- c=0.15 → 17 anomalies
- c=0.20 → 21 anomalies""",
        FIGS / "fig_top6_multipanel.png")

    _add_content_slide(prs,
        "발견 ④ Heckman 2-stage 패널 회귀",
        """**Stage 2 결과 (N=104, 23 firms)**

| 변수 | β | Bootstrap 95% CI |
|---|---|---|
| **ln(GIR)** | **−2.00** | **[−8.93, −0.21]** ⭐ |
| in_KSSB_30 | (perfect predict) | dropped |
| IMR | −12.69 | [−46.70, +2.26] |
| yr_2020 | −2.36 | [−12.26, +1.13] |
| **yr_2021** | **+11.49** | **[+0.13, +32.98]** ⭐ |
| yr_2022 | −6.60 | [−29.10, +2.82] |
| yr_2023 | −6.30 | [−28.44, +3.05] |

**해석**
- 대기업일수록 괴리율 낮음 (정확 보고)
- 2021년 일시적 괴리 급증 (코로나 회복)
- IMR borderline (선택편향 존재 시사)
- 산업 더미 통제 + cluster SE by firm""")

    _add_content_slide(prs,
        "SHAP 기여도 분해 (전체 firms)",
        """**SHAP Feature Importance**
mean |SHAP| 순위:
1. NO₂_resid (기상보정 잔차)
2. log(GIR Scope 1)
3. 괴리율 DIFF_rel
4. SO₂_resid
5. ODIAC CO₂
6. CO_resid
7. ERA5 BLH
8. energy_TJ

**Implication**
- 위성 신호가 GIR 절대값보다 더 강한 예측력
- 4중 비교의 가치 입증
- ODIAC이 지지하는 결정적 firms 식별

**TreeExplainer 옵션**
feature_perturbation='interventional'
→ Path-dependent 편향 회피""",
        FIGS / "fig_shap_summary.png")

    _add_content_slide(prs,
        "SHAP Waterfall — Top 5 이상 firms",
        """**개별 firm 분해 (Top 5)**

각 firm의 이상 점수가
어떤 feature 조합으로 결정됐는지 분해

**한국전력공사**
- log(GIR) 절대 규모가 압도적 기여
- 위성 신호 중간

**포스코홀딩스 (3년 연속 transient)**
- NO₂_resid 하락 + GIR 상승 방향 불일치
- ODIAC CO₂도 동일 방향 → 강한 신호

**Interpretation**
- Firm-specific 정책 권고 가능
- "왜 이 firm이 이상으로 탐지됐는가"
  를 인간이 해석 가능한 언어로 제시""",
        FIGS / "fig_shap_waterfall_top5.png")

    _add_content_slide(prs,
        "검증 우선순위 매트릭스 (Top 10)",
        """**priority_score 계산식**
0.4·괴리심각도 + 0.4·위성불일치 + 0.2·이상등급

**상위 10개사 (KSSB 2028 대상)**
1. 한국전력공사     0.54
2. 포스코홀딩스     0.47
3. LG에너지솔루션   0.40
4. 네이버           0.40
5. CJ제일제당       0.38
6. LG디스플레이     0.31
7. SK하이닉스       0.30
8. 이마트           0.30
9. 롯데쇼핑         0.30
10. 삼성전자        0.28 (D 패턴)

**활용**
- 상위 25% → 즉시 검증 대상 지정
- KEITI 환경부 검증 자원 차등 배분
- KSSB 시행령에 priority 명시""",
        FIGS / "fig_priority_matrix.png")

    # ============ DECK 5: 정책 + 결론 (41-50) ============
    _add_title_slide(prs,
        "Part 5 — 정책 함의와 결론",
        "3 정책 카드 + KSSB 활용 + 한계 + 향후")

    _add_content_slide(prs,
        "정책 카드 ① — KEITI 환경책임투자 플랫폼 DRI",
        """**Disclosure Reliability Index (DRI) 신설**

KEITI 환경책임투자 플랫폼에
4중 검증 신뢰성 지수 편입

**DRI 구성요소**
- GIR-ESG 괴리율 (40%)
- 위성·ODIAC 일관성 점수 (40%)
- 이상탐지 등급 (20%)

**활용 방식**
- 투자자 공개 지수
- ESG 펀드 운용 reference
- 책임 투자 자본 배분 정확성 향상

**선행 사례**
- MSCI ESG Rating + Sustainalytics
- 본 연구는 위성 기반 객관 측정 추가""")

    _add_content_slide(prs,
        "정책 카드 ② — 우선순위 매트릭스 차등 검증",
        """**KSSB 2028 49개사 차등 자원 배분**

상위 25% → 즉시 검증 대상
  - 환경부 + 검증기관 합동 현장 점검
  - 패턴 D · 구조적 이상 firms 우선

26-50% → 우선 관찰
  - 분기별 데이터 모니터링
  - 위성 신호 trend 분석

51-100% → 일반 모니터링
  - 연 1회 routine 점검

**기대 효과**
- 한정된 검증 자원의 효율 배분
- 무작위 선별 대비 검증 효율 극대화
- 데이터 기반 정책 결정 정착""",
        FIGS / "fig_priority_matrix.png")

    _add_content_slide(prs,
        "정책 카드 ③ — KSSB 제2호 시행령 강화",
        """**3단계 프로토콜 의무화**

Step 1: GIR-ESG 대조표 첨부
  - 모든 의무공시 보고서에
    GIR vs ESG 괴리율 표 의무 첨부
  - 괴리율 ±20% 초과 시 사유 설명 의무

Step 2: 위성 모니터링 연계
  - Sentinel-5P 사업장 buffer 신호
  - ODIAC CO₂ 1km 비교
  - 환경부 → KSSB 정기 보고

Step 3: 독립 검증기관 지정
  - 패턴 D 또는 괴리율 ±20% 초과 firms
  - 환경부 지정 검증기관 현장 검증
  - EU CBAM 대응 탄소회계 신뢰성 확보""")

    _add_content_slide(prs,
        "KSSB 2028 즉시 활용 경로",
        """**Gold 23사 → KSSB 49사 매핑**

본 연구 대상이 KSSB 1차 49개사와
대부분 교집합 → 즉시 적용 가능

**3가지 즉시 활용 경로**
1. 검증기관 인력 양성
   - 위성 분석 역량 인증 프로그램
   - 본 연구 도구 (run_all_analysis.py)
     교육 자료로 활용

2. 시행령 제정
   - 본 연구 priority_score 직접 인용
   - 패턴 D 정의 → 우선 검증 기준

3. GEE 시스템 구축
   - 환경부·KSSB 공동 GEE 워크스페이스
   - 자동화 수집·파싱 파이프라인 적용""")

    _add_content_slide(prs,
        "종합 논의 — 패턴 D 가설 3가지",
        """**가설 1: 보고경계 확대**
- 해외 자회사 신규 통합
- 연결 범위 변경
→ ESG 절대값 증가, 국내 GIR과 분리

**가설 2: Scope 정의 변경**
- Scope 1 → Scope 2 reclassification
- 자가발전 vs 외부 전력 처리 차이
→ 분류 변경에 따른 절대값 변화

**가설 3: 실제 배출 감소 + 보고 시차**
- 효율 개선 (CO₂ 회수, 폐열 활용)
- 보고는 작년 데이터 → 1년 시차
→ 위성은 즉시 반영, 보고는 지연

**검증 방법**
본 분석으로 단정 불가 — 후속 현장 검증 필요
인과 추론 시도하지 않고 기술적 서술만 유지""")

    _add_content_slide(prs,
        "방법론 강건성 종합",
        """**ERA5 보정 강건성**
- 보정 전후 패턴 분류 동일
- HCHO R²=0.94, NO₂ R²=0.76
- MERRA-2 sensitivity 일치

**이상탐지 강건성**
- contamination sweep (0.05-0.20)
- IF + LOF 앙상블 false positive 감소
- KCGS 21건 supervised 검증

**회귀 강건성**
- Bootstrap 95% CI (B=2000, firm-block)
- cluster SE by firm
- industry × year fixed effects

**ASOS 지상 검증**
- ERA5 vs ASOS 5 stations
- 평균 거리 9.2 km
- 격자 모델 vs 지점 관측 sanity check""")

    _add_content_slide(prs,
        "한계 및 대응 매트릭스 (8가지)",
        """**L1 소샘플** (N=23) → Bootstrap CI, Silver 확장
**L2 인과 추론 불가** → "관찰" 기술적 서술만
**L3 조직경계 잔존 불확실성** → 3계층 샘플
**L4 NO₂ ≠ CO₂ 등가 아님** → ODIAC 추가, 5채널 과반
**L5 본사 좌표 사용 firms** → 8.0 신규 한계 명시
**L6 5년 시계열 짧음** → MK p<0.1 + |τ|≥0.4 완화
**L7 contamination 자의성** → KCGS 라벨 calibration
**L8 ESG PDF 11 구조적 미발간** → 명문화 (네이버 2019 등)

**핵심: 본 연구의 한계 transparency**
모든 한계가 보고서·코드·ADR에 명시되어
후속 연구가 즉시 boundary 인지 가능""")

    _add_content_slide(prs,
        "향후 과제",
        """**확장 1: Silver 205사로 통계적 검정력 강화**
- 패널 N=23 → 1,025 firm-year
- 회귀 자유도 충분 → 강건한 인과 추론

**확장 2: LSTM Autoencoder 시계열 이상탐지**
- 현재 Mann-Kendall (단순)
- LSTM-AE → 비선형 패턴 + 장기 의존성

**확장 3: Sentinel-5P CH₄ 추가**
- 에너지 업종 메탄 배출 검증
- LNG·정유·축산 firms 확장

**확장 4: 다중 buffer 가중합**
- 사업장 복수 보유 firms
- 사업장별 조업률 가중
- 더 정확한 firm-level 신호

**확장 5: KSSB 2028 시행 후 비교**
- 의무공시 시행 전후 패턴 변화
- 본 연구 검증 체계의 효과성 평가""")

    _add_content_slide(prs,
        "핵심 기여 한 문장 + Q&A",
        """**한 줄 정리**

"한국 코스피 상장기업의 GIR 법정 배출량,
ESG 자체보고, Sentinel-5P 위성 4종, ODIAC CO₂를
ERA5 기상보정 후 4중으로 비교해
Mann-Kendall 기반 공시 불일치 5종 패턴을 분류하고,
Heckman·이상탐지 3층·SHAP를 통해 KSSB 2028
의무공시 검증 우선순위 프레임워크를 설계하여
포스코홀딩스(τ=±1.00 극단)와 삼성전자에서
패턴 D를 처음으로 정량적으로 식별했다."

**Q&A 컨택트**
zxsa0716@kookmin.ac.kr
GitHub: https://github.com/zxsa0716/AX_Contest

**감사합니다**""",
        FIGS / "fig_concept_4channel.png")

    # ============ DECK 6: 23개사 심층 분석 (10) ============
    _add_title_slide(prs,
        "Part 6 — 23개사 산업별 심층 분석",
        "패턴 분류 결과의 사업 맥락 해석 + 8 신규 시각화")

    _add_content_slide(prs,
        "Sentinel-5P 4종 위성 신호 한국 분포 (2023)",
        """**4 panel 비교**
- NO₂ (위, 좌): 산업·교통 NOx 배출 추적자
- SO₂ (위, 우): 화력발전·제철 특화
- CO (아래, 좌): 불완전 연소 (장수명)
- HCHO (아래, 우): 석유화학 VOC

**관찰**
- NO₂: 수도권·울산·여수·포항 hotspot
- SO₂: 한정된 석탄·제철 cluster
- CO: 광범위 (수송 + 산업)
- HCHO: 여수·울산 석유화학 단지 집중

**해석**
4 species 동시 비교로 산업 특성에 맞는
위성 채널 선택 가능 (예: 철강은 NO₂+SO₂)""",
        FIGS / "fig_s5p_4species_korea.png")

    _add_content_slide(prs,
        "ODIAC CO₂ 2019 vs 2023 변화 (한국 전체)",
        """**3 panel 비교**
- 좌: 2019-07 ODIAC
- 중: 2023-07 ODIAC
- 우: 차이 (2023-2019, 빨강=증가)

**관찰**
- 수도권 일부 영역 증가 (데이터센터·인구)
- 일부 산업단지 감소 (탈탄소 효과)
- 전국 합산은 약간 감소 추세

**시사점**
ODIAC 1km 해상도로 firm-level 변화 가시화
패턴 D firms (포스코·삼성전자)의 하강 추세
ODIAC 데이터로 일관 확인됨""",
        FIGS / "fig_odiac_change_2019_2023.png")

    _add_content_slide(prs,
        "업종별 GIR 5년 추이 (6 산업군)",
        """**6 panel: 업종별 firm 시계열**
- steel (POSCO, 현대제철)
- petrochem (롯데케미칼, 한화솔루션, LG엔솔)
- power_coal (KEPCO, SK이노)
- semiconductor (삼성전자, SK하이닉스, LGD)
- finance (네이버, KT, 삼성생명, IBK)
- other (나머지)

**산업별 패턴 차이**
- steel: 절대값 큰 변동 (수출 시황)
- petrochem: 일반적 감소 추세
- finance: Scope 1 미미, 안정
- semicon: 사업 확장 + 효율 trade-off""",
        FIGS / "fig_industry_timeseries.png")

    _add_content_slide(prs,
        "23개사 4채널 정규화 추이 (2019=100)",
        """**Small multiples grid (5×5)**
- 모든 23 firms 한 화면 비교
- GIR(검정)·ESG(회색)·NO₂(주황)·ODIAC(파랑)
- 2019=100으로 정규화 → 변화율 비교

**시각적 발견**
- 포스코홀딩스: GIR↑↑ vs 위성↓↓ 명확
- 삼성전자: ESG가 가장 가파른 상승
- 네이버: 4채널 모두 상승 (확장)
- 대부분 firms: 위성·ODIAC 동반 하락
- 금융 firms: GIR/ESG 변동 미미

→ 산업별 다양성 + 공통 추세 동시 가시화""",
        FIGS / "fig_all23_normalized.png")

    _add_content_slide(prs,
        "포스코홀딩스 4중 비교 디테일",
        """**4 panel (각 trend line + R²)**
- GIR 법정: τ=+1.00 완벽 상승
- ESG 자체: τ=+0.67 상승
- 위성 NO₂: τ=-1.00 완벽 하강
- ODIAC CO₂: τ=-1.00 완벽 하강

**해석 (단정 회피)**
가설 1: 보고경계 확장 (해외 자회사)
가설 2: Scope 정의 변경 reclassification
가설 3: 효율 개선 + 보고 시차

**KSSB 2028 권고**
포스코홀딩스 우선 검증 1순위 지정
환경부 + KEITI 합동 현장 점검
GIR-ESG 교차 검증 + 위성 데이터 활용""",
        FIGS / "fig_posco_4channel_detail.png")

    _add_content_slide(prs,
        "Heckman 회귀 계수 Forest Plot",
        """**Bootstrap 95% CI 시각화**
- 주황: 유의 (CI가 0 미포함)
- 회색: 비유의

**유의한 계수**
- ln(GIR): -2.00 [-8.93, -0.21]
- yr_2021: +11.49 [+0.13, +32.98]

**비유의 그러나 시사점 있는 계수**
- IMR: -12.69 (선택편향 borderline)
- 산업·연도 dummies (대부분)

**해석**
대기업일수록 괴리율 낮음
2021년 코로나 회복기 일시적 괴리 급증
산업별 차이는 sample 한계로 강한 결론 X""",
        FIGS / "fig_heckman_forest.png")

    _add_content_slide(prs,
        "이상탐지 2D 분포 (괴리 × 위성 불일치)",
        """**X축**: GIR-ESG 괴리도
**Y축**: GIR-위성·ODIAC 불일치도

**우상단: 즉시 검증 대상**
- 포스코홀딩스 (D 패턴, 큰 점)
- 삼성전자 (D 패턴)
- 현대모비스 (C 패턴)

**좌하단: 정상 권역**
- A_consistent_down 12개사 다수
- 금융·서비스 firms

**우하단: 위성만 불일치**
- 일부 산업 firms (본사 좌표 한계)

**좌상단: 공시 채널만 불일치**
- mixed 패턴 firms (부분 일치)""",
        FIGS / "fig_anomaly_2d.png")

    _add_content_slide(prs,
        "업종별 괴리율 분포 (Box Plot)",
        """**6 산업군 GIR-ESG 괴리율**

(ESG − GIR) / GIR × 100 (%)

**분포 특징**
- steel: 중앙값 0 근처, 변동 작음
- petrochem: 양의 편향 (ESG 더 큼)
- finance: 큰 변동 (Scope 1 작아 비율 큼)
- semicon: 양의 편향 (사업 확장)
- power_coal: 음의 편향 (KEPCO 효과)
- other: 가장 큰 변동

**시사점**
산업 특성별 공시 패턴 차이 확인
KSSB 의무공시 시행 시 산업별
검증 가이드라인 차별화 필요""",
        FIGS / "fig_industry_boxplot.png")

    _add_content_slide(prs,
        "Part 6 결론: 4가지 산업별 인사이트",
        """**1. D 패턴은 철강+반도체 두 산업에만**
   포스코·삼성전자 — 절대 배출 큼 + 사업구조 변화 활발

**2. C 패턴은 현대모비스 단 1건**
   ESG 보고 범위 확대 시점 변동성 가설

**3. 위성 NO₂ 12개 firm 단조 하강**
   한국 정부 NOx 감축 정책 (2017-2022) 효과
   산업별 차이 < 정책 공통 효과

**4. 금융·서비스 firms (네이버 제외) 경계 사례**
   Scope 1 미미 + 도시 배경 매몰
   → Scope 2·3 검증이 더 중요한 의제

**핵심**: 본 4중 비교 프레임은 산업·중공업
firms에서 가장 강력한 진단력 제공""")

    prs.save(OUT)
    print(f"[saved] {OUT}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"File size: {OUT.stat().st_size:,} bytes")


if __name__ == "__main__":
    build()
